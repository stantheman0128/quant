# -*- coding: utf-8 -*-
"""Build KOL poster v4 (ENGLISH body) cloning the format of '12_3D Anomaly detection.pdf':
   19.05 x 27.52 cm, all Times serif, red section bars (white bold serif title),
   flat colored flow boxes + arrows, two columns, LaTeX-booktabs tables.
   Logo: uses user-supplied 'logo.png' if present, else '_logo.png' (extracted from 海報範例.pptx)."""
import os
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette ----------
RED   = "9A2A2A"
INK   = "1A1A2D"; BLACK="000000"
NAVY  = "1A1A2D"; BLUE="1F78B4"; GREEN="1A9541"; TEAL="0E8F6E"; PURPLE="8A2E9F"; DRED="9A2A2A"
LAT, EA = "Times New Roman", "PMingLiU"

LOGO = "logo.png" if os.path.exists("logo.png") else ("logo.jpg" if os.path.exists("logo.jpg") else "_logo.png")

prs = Presentation()
prs.slide_width  = Cm(19.05); prs.slide_height = Cm(27.52)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def set_run(r, size, color=BLACK, bold=False, italic=False, latin=LAT, ea=EA):
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = RGBColor.from_string(color); r.font.name = latin
    rPr = r._r.get_or_add_rPr()
    for tag, face in (('a:ea', ea), ('a:cs', latin)):
        e = rPr.find(qn(tag))
        if e is None: e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set('typeface', face)

def box(x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = RGBColor.from_string(line); sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def textbox(x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            line_spacing=1.0, space_after=2, wrap=True):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = Cm(0.1); tf.margin_right = Cm(0.1); tf.margin_top = Cm(0.03); tf.margin_bottom = Cm(0.03)
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = align; p.line_spacing = line_spacing
        p.space_before = Pt(0); p.space_after = Pt(space_after)
        for text, kw in para:
            r = p.add_run(); r.text = text; set_run(r, **kw)
    return tb

def bar(x, y, w, en, h=0.72, size=15):
    box(x, y, w, h, fill=RED)
    tb = slide.shapes.add_textbox(Cm(x+0.18), Cm(y), Cm(w-0.3), Cm(h))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    tf.margin_top = Cm(0.02); tf.margin_bottom = Cm(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = en; set_run(r, size, "FFFFFF", bold=True)
    return y + h

def rule(x, y, w, weight=1.4):
    ln = slide.shapes.add_connector(2, Cm(x), Cm(y), Cm(x+w), Cm(y))
    ln.line.color.rgb = RGBColor.from_string(BLACK); ln.line.width = Pt(weight)
    return ln

# ============================ HEADER ============================
slide.shapes.add_picture(LOGO, Cm(0.49), Cm(0.42), Cm(2.3), Cm(2.3))
nb = box(17.05, 0.42, 1.5, 1.5, fill="FFFFFF", line=BLACK, line_w=1.25)
tf = nb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "50"; set_run(r, 22, BLACK)
# title (English main)
textbox(2.7, 0.55, 14.2, 2.0, [
    [("Chinese Financial KOL Opinions as a Stock-Selection Factor:", dict(size=17, color=BLACK, bold=True))],
    [("Testing Incremental Information over Alpha101", dict(size=17, color=BLACK, bold=True))],
], align=PP_ALIGN.CENTER, line_spacing=1.04, space_after=2)
textbox(3.0, 2.78, 9.0, 0.6, [[("Po-Han Shih (施博瀚)", dict(size=12, color=BLACK))]], align=PP_ALIGN.CENTER)
textbox(11.6, 2.82, 6.9, 0.6, [[("Advisor: Prof. Po-Wen Chi (紀博文)", dict(size=10.5, color=BLACK))]], align=PP_ALIGN.RIGHT)
rule(0.0, 3.48, 19.05, 1.6)

# ============================ FLOW DIAGRAM ============================
flow = [("KOL Posts","FB",NAVY), ("LLM-as-Formatter","format, no scoring",BLUE),
        ("4-Field Schema","{t,id,dir,w}",GREEN), ("Daily KOL Factor","f_kol",TEAL),
        ("Orthogonalize","remove Alpha101",PURPLE), ("Incremental IC","residual signal",DRED)]
fx, fy, fw, fh, gap = 0.5, 3.8, 2.72, 1.55, 0.34
for i,(en,sub,col) in enumerate(flow):
    x = fx + i*(fw+gap)
    box(x, fy, fw, fh, fill=col)
    tb = slide.shapes.add_textbox(Cm(x+0.05), Cm(fy), Cm(fw-0.1), Cm(fh))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    tf.margin_left=Cm(0.04); tf.margin_right=Cm(0.04); tf.margin_top=Cm(0.02); tf.margin_bottom=Cm(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing=1.0
    r = p.add_run(); r.text = en; set_run(r, 9.5, "FFFFFF", bold=True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before=Pt(1)
    r2 = p2.add_run(); r2.text = sub; set_run(r2, 7.5, "EAE0E0", italic=True)
    if i < len(flow)-1:
        box(x+fw+0.02, fy+fh/2-0.18, gap-0.04, 0.36, fill="555555", shape=MSO_SHAPE.RIGHT_ARROW)
textbox(0.5, fy+fh+0.12, 18.05, 0.55, [[
    ("weight", dict(size=8.5, color=RED, bold=True, latin="Consolas", ea="Consolas")),
    (" = the KOL's walk-forward hit-rate on posts BEFORE this one (point-in-time, no look-ahead); direction is extracted by the LLM.",
     dict(size=9, color="444444"))]], align=PP_ALIGN.CENTER)

# ============================ COLUMNS ============================
LX, RX, CW = 0.5, 9.75, 8.8
TOP = 6.85

# ---- LEFT: Motivation ----
y = bar(LX, TOP, CW, "Motivation")
textbox(LX, y+0.12, CW, 5.4, [
  [("Retail traders make up ", dict(size=10)), ("~50%", dict(size=10,bold=True)),
   (" of TWSE turnover, and many follow financial KOLs. Prior work covers print advisories, sell-side analysts and TV pundits, but ", dict(size=10)),
   ("independent, unregulated online financial KOLs are under-explored", dict(size=10,bold=True,color=RED)),
   (" — and, free of client / regulatory conflicts, form a relatively clean sample.", dict(size=10))],
  [("KOL posts are free text. Rather than quantify a KOL's causal ", dict(size=10)),
   ("“influence”", dict(size=10,italic=True)),
   (", we treat KOL opinion as a ", dict(size=10)),
   ("stock-selection factor", dict(size=10,bold=True,color=RED)),
   (" and ask a testable question:", dict(size=10))],
  [("After removing known price-volume factors (Alpha101), does the KOL factor still predict returns?", dict(size=10,bold=True)),
   ("  = new information, or just momentum repackaged?", dict(size=10,italic=True,color="444444"))],
], align=PP_ALIGN.LEFT, line_spacing=1.07, space_after=4)

# ---- LEFT: Methodology ----
y = bar(LX, 13.75, CW, "Methodology")
textbox(LX, y+0.12, CW, 10.6, [
  [("① LLM-as-Formatter (no scoring): ", dict(size=10.5,bold=True,color=BLUE)),
   ("free text → 4-field schema. The LLM only formats and disambiguates entities (台積 ≡ 2330 ≡ TSMC); it outputs no confidence score.", dict(size=10))],
  [("② Walk-forward weight (point-in-time): ", dict(size=10.5,bold=True,color=GREEN)),
   ("weight = the KOL's hit-rate on posts BEFORE this one; estimated only with ≥ 30 prior signals, else 0.5. Strictly prevents look-ahead.", dict(size=10))],
  [("        weight = hits / n   (only prior signals)", dict(size=9.5,latin="Consolas",ea="Consolas",color=INK))],
  [("③ Aggregate into a daily KOL factor: ", dict(size=10.5,bold=True,color=INK)),
   ("pool recent KOL signals per stock → cross-sectional factor f_kol.", dict(size=10))],
  [("④ Incremental IC vs Alpha101 (core): ", dict(size=10.5,bold=True,color=DRED)),
   ("cross-sectionally regress f_kol on Alpha101, take the residual e (what Alpha101 cannot explain), and correlate e with forward returns.", dict(size=10))],
  [("        IC_incr = corr( resid(f_kol | Alpha101), ret )", dict(size=9.5,latin="Consolas",ea="Consolas",color=INK))],
  [("→ remains → new information;  vanishes → just momentum (an honest null result).", dict(size=10,italic=True,color="444444"))],
], align=PP_ALIGN.LEFT, line_spacing=1.08, space_after=4)

# ---- LEFT: Reference ----
y = bar(LX, 24.5, CW, "Reference", h=0.62, size=13)
textbox(LX, y+0.06, CW, 2.6, [
  [("[1] Kakushadze (2016). 101 Formulaic Alphas. Wilmott.", dict(size=8.5,color="333333"))],
  [("[2] Pardo (2008). The Evaluation & Optimization of Trading Strategies.", dict(size=8.5,color="333333"))],
  [("[3] Bailey & López de Prado (2014). The Deflated Sharpe Ratio.", dict(size=8.5,color="333333"))],
  [("[4] Kakhbod et al. (2023). Finfluencers.", dict(size=8.5,color="333333"))],
], align=PP_ALIGN.LEFT, line_spacing=1.0, space_after=1)

# ---- RIGHT: Experimental Results ----
y = bar(RX, TOP, CW, "Experimental Results")
textbox(RX, y+0.1, CW, 0.55, [[
  ("PoC: 3 FB KOLs · ", dict(size=10,color="444444")), ("1,685", dict(size=10,bold=True)),
  (" opinions · ", dict(size=10,color="444444")), ("404", dict(size=10,bold=True)),
  (" tickers · 2024–2026", dict(size=10,color="444444"))]], align=PP_ALIGN.LEFT)

def booktabs(x, y, w, headers, rows, col_w, ours_idx, row_h=0.62, hsz=10.5, bsz=10):
    n = len(rows)+1; total_h = row_h*n
    rule(x, y, w, 1.6)
    cx = x
    for j,htext in enumerate(headers):
        tb = slide.shapes.add_textbox(Cm(cx), Cm(y+0.02), Cm(col_w[j]), Cm(row_h))
        tf=tb.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; tf.word_wrap=True
        tf.margin_left=Cm(0.06); tf.margin_right=Cm(0.04); tf.margin_top=Cm(0); tf.margin_bottom=Cm(0)
        p=tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
        r=p.add_run(); r.text=htext; set_run(r, hsz, BLACK, bold=True)
        cx += col_w[j]
    rule(x, y+row_h, w, 1.0)
    for i,row in enumerate(rows):
        ry = y + row_h*(i+1); is_ours = (i==ours_idx); cx = x
        for j,cell in enumerate(row):
            tb = slide.shapes.add_textbox(Cm(cx), Cm(ry+0.02), Cm(col_w[j]), Cm(row_h))
            tf=tb.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; tf.word_wrap=True
            tf.margin_left=Cm(0.06); tf.margin_right=Cm(0.04); tf.margin_top=Cm(0); tf.margin_bottom=Cm(0)
            p=tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
            r=p.add_run(); r.text=cell; set_run(r, bsz, DRED if is_ours else BLACK, bold=is_ours)
            cx += col_w[j]
    rule(x, y+total_h, w, 1.6)
    return y+total_h

ey = booktabs(RX, y+0.75, CW, ["Method","hit","IC","Sharpe"],
    [["Direction-only","54.3%","+0.045","+1.11"],
     ["Walk-forward weighted","55.0%","+0.097","+0.97"],
     ["High-confidence filter (Ours)","57.9%","—","+1.65"]],
    col_w=[4.6,1.4,1.5,1.3], ours_idx=2)
textbox(RX, ey+0.06, CW, 0.5, [[("Table 1. ", dict(size=8.5,bold=True)),
    ("Three signal strategies; quality beats quantity.", dict(size=8.5,color="444444"))]], align=PP_ALIGN.CENTER)

ey2 = booktabs(RX, ey+0.78, CW, ["weight quantile","n","hit rate"],
    [["< 0.45","276","49.3%"],
     ["0.55 – 0.60","271","55.0%"],
     ["> 0.60 (Ours)","180","62.2%"]],
    col_w=[4.4,2.2,2.2], ours_idx=2)
textbox(RX, ey2+0.06, CW, 0.5, [[("Table 2. ", dict(size=8.5,bold=True)),
    ("Hit-rate by weight quantile: 49% → 62%, monotonic.", dict(size=8.5,color="444444"))]], align=PP_ALIGN.CENTER)
textbox(RX, ey2+0.62, CW, 0.5, [[("⚠ Raw IC, not yet residualized on Alpha101; incremental IC (core check) in progress.",
    dict(size=9,color="B5651D",italic=True))]], align=PP_ALIGN.LEFT)

# ---- RIGHT: Our Ideas ----
y = bar(RX, ey2+1.35, CW, "Our Ideas")
textbox(RX, y+0.12, CW, 6.0, [
  [("Turn “is the KOL just momentum?” into a falsifiable number. ", dict(size=10.5,bold=True)),
   ("Incremental IC vs Alpha101 directly measures the KOL factor's increment over known price-volume factors (target ≥ 0.03 / ≥ 0.05).", dict(size=10.5))],
  [("Separation of duties + no look-ahead. ", dict(size=10.5,bold=True)),
   ("The LLM only formats; scoring is externalized to historical walk-forward. Treating KOL as a factor (not “influence”) keeps conclusions strictly back-testable.", dict(size=10.5))],
  [("Positioning (honest). ", dict(size=10.5,bold=True)),
   ("The methods are standard quant; the novelty is the sample — independent Traditional-Chinese KOLs — and bringing this rigor to an under-explored area.", dict(size=10.5))],
  [("Signal structure already emerging. ", dict(size=10.5,bold=True)),
   ("The PoC surfaced 4 anomalies: KOL skill does not transfer across markets; overconfidence underperforms; an alleged “contra-indicator” is really a lagged follower; long/short skill is asymmetric.", dict(size=10.5))],
], align=PP_ALIGN.LEFT, line_spacing=1.12, space_after=7)

out = "KOL_Poster_v4.pptx"
try:
    prs.save(out)
except PermissionError:
    out = "_v4_build.pptx"; prs.save(out)   # v4 locked (open in PowerPoint) -> temp
print("saved", out, " (logo:", LOGO + ")")
